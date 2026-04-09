#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
六维耦合系统 PPT 生成脚本 - 小爪 🐾
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os

# 创建 PPT
prs = Presentation()

# 设置宽屏 (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
PRIMARY_COLOR = RGBColor(0, 51, 102)      # 深蓝
ACCENT_COLOR = RGBColor(0, 153, 204)      # 浅蓝
HIGHLIGHT_COLOR = RGBColor(255, 102, 0)   # 橙色
TEXT_COLOR = RGBColor(50, 50, 50)         # 深灰

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    return slide

def add_content_slide(prs, title, content_items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 设置内容
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(14)
        
        if item.startswith("•"):
            p.level = 1
    
    return slide

def add_diagram_slide(prs, title, description):
    """添加图示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白版式
    
    # 设置标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 添加中心图形（六维示意）
    left = Inches(4)
    top = Inches(2)
    width = Inches(5.333)
    height = Inches(4)
    
    # 添加六边形
    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.fill.transparency = 0.7
    shape.line.color.rgb = PRIMARY_COLOR
    shape.line.width = Pt(3)
    
    # 添加说明文字
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(18)
    desc_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 开始创建幻灯片 ============

# 1. 标题页
add_title_slide(
    prs,
    "六维耦合系统",
    "Six-Degree-of-Freedom Coupled System\n技术原理与应用\n" + datetime.now().strftime("%Y年%m月")
)

# 2. 目录
add_content_slide(
    prs,
    "目录",
    [
        "1. 六维系统概述",
        "2. 耦合系统原理",
        "3. 数学模型",
        "4. 关键技术",
        "5. 应用领域",
        "6. 发展趋势"
    ]
)

# 3. 六维系统概述
add_content_slide(
    prs,
    "六维系统概述",
    [
        "什么是六维系统？",
        "• 六自由度（6-DOF）运动系统",
        "• 三个平移自由度：X、Y、Z 轴",
        "• 三个旋转自由度：Roll、Pitch、Yaw",
        "",
        "系统特点：",
        "• 全空间运动能力",
        "• 高精度定位控制",
        "• 多轴协同耦合"
    ]
)

# 4. 六维示意图
slide = add_diagram_slide(
    prs,
    "六自由度示意图",
    "六个自由度：X/Y/Z 平移 + Roll/Pitch/Yaw 旋转"
)

# 5. 耦合系统原理
add_content_slide(
    prs,
    "耦合系统原理",
    [
        "耦合的定义",
        "• 多个子系统相互影响、相互作用",
        "• 输入输出之间存在关联关系",
        "",
        "耦合类型：",
        "• 机械耦合：结构连接导致的力/力矩传递",
        "• 控制耦合：多控制器之间的信号交互",
        "• 动力耦合：能量在系统间的传递",
        "",
        "解耦控制：",
        "• 将多变量系统分解为单变量系统",
        "• 降低控制复杂度"
    ]
)

# 6. 数学模型
add_content_slide(
    prs,
    "数学模型",
    [
        "运动学方程：",
        "• 位置向量：p = [x, y, z]ᵀ",
        "• 姿态向量：θ = [α, β, γ]ᵀ",
        "• 广义坐标：q = [pᵀ, θᵀ]ᵀ",
        "",
        "动力学方程：",
        "• M(q)q̈ + C(q,q̇)q̇ + G(q) = τ",
        "• M: 质量矩阵",
        "• C: 科里奥利力矩阵",
        "• G: 重力向量",
        "• τ: 广义力/力矩"
    ]
)

# 7. 关键技术
add_content_slide(
    prs,
    "关键技术",
    [
        "1. 精密驱动技术",
        "• 伺服电机/直线电机",
        "• 压电陶瓷驱动器",
        "",
        "2. 传感检测技术",
        "• 激光干涉仪",
        "• 光电编码器",
        "• 惯性测量单元 (IMU)",
        "",
        "3. 控制算法",
        "• PID 控制",
        "• 自适应控制",
        "• 鲁棒控制",
        "• 智能控制"
    ]
)

# 8. 应用领域
add_content_slide(
    prs,
    "应用领域",
    [
        "🔬 科学实验",
        "• 振动模拟测试台",
        "• 微重力环境模拟",
        "",
        "🏭 工业制造",
        "• 精密加工定位",
        "• 机器人末端执行器",
        "",
        "✈️ 航空航天",
        "• 飞行模拟器",
        "• 卫星姿态控制",
        "",
        "🏥 医疗设备",
        "• 手术机器人",
        "• 康复训练设备"
    ]
)

# 9. 发展趋势
add_content_slide(
    prs,
    "发展趋势",
    [
        "1. 高精度化",
        "• 纳米级定位精度",
        "• 亚微弧度角精度",
        "",
        "2. 高速化",
        "• 高频响驱动",
        "• 实时控制优化",
        "",
        "3. 智能化",
        "• AI 辅助控制",
        "• 自学习自适应",
        "",
        "4. 集成化",
        "• 模块化设计",
        "• 系统小型化"
    ]
)

# 10. 总结
add_content_slide(
    prs,
    "总结",
    [
        "六维耦合系统核心价值：",
        "• 实现全空间六自由度精确控制",
        "• 解决多轴协同耦合难题",
        "• 支撑高端装备核心技术",
        "",
        "技术挑战：",
        "• 强耦合非线性",
        "• 高精度要求",
        "• 实时性约束",
        "",
        "未来展望：",
        "• 更智能、更精密、更集成"
    ]
)

# 11. 结束页
add_title_slide(
    prs,
    "谢谢观看",
    "Q & A\n\n小爪 🐾 自动生成"
)

# 保存 PPT
output_path = "/home/tenbox/Desktop/六维耦合系统.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
