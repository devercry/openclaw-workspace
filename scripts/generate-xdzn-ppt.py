#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
芯动智能科技（温州）有限公司介绍 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_xdzn_ppt():
    prs = Presentation()
    
    # ==================== 第 1 页：封面 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "芯动智能科技（温州）有限公司"
    subtitle.text = "智能制造 · 一站式集成\n精密滑台 · 工业数字化 · 项目交付"
    
    # ==================== 第 2 页：公司简介 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "公司简介"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "北大 - 温州激光与光电子联合研发中心孵化企业"
    
    p = tf.add_paragraph()
    p.text = "专注于纳米级精密滑台及光机组件的研发与生产，达到国际先进水平"
    
    p = tf.add_paragraph()
    p.text = "产品兼具高刚性、低噪音与优异的速度稳定性"
    
    p = tf.add_paragraph()
    p.text = "支持多轴定制与嵌入式控制系统集成"
    
    p = tf.add_paragraph()
    p.text = "注册地：浙江温州"
    
    p = tf.add_paragraph()
    p.text = "业务聚焦：智能装备 · 工业数字化"
    
    # ==================== 第 3 页：企业文化 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "企业文化"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # 企业定位
    p = tf.add_paragraph()
    p.text = "🎯 企业定位"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "全球精密滑台领域的创新驱动型企业，专注于高精度、高可靠性的运动控制解决方案"
    p.level = 1
    
    # 企业使命
    p = tf.add_paragraph()
    p.text = "🚀 企业使命"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "推动精密运动控制技术的创新与应用，为全球工业智能化和高质量发展提供卓越解决方案"
    p.level = 1
    
    # 企业愿景
    p = tf.add_paragraph()
    p.text = "🌟 企业愿景"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "成为全球精密滑台领域的领导者，通过持续创新和卓越品质，为客户提供高效、精准、可靠的解决方案"
    p.level = 1
    
    # 服务理念
    p = tf.add_paragraph()
    p.text = "💡 服务理念"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "专业、高效、诚信、贴心 — 以客户为中心，提供全方位、定制化的解决方案"
    p.level = 1
    
    # ==================== 第 4 页：核心产品 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "核心产品"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "纳米级精密滑台"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "高精度、高刚性、低噪音的运动控制核心部件"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "光机组件"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "光学平台、线性滑台、镜筒等精密光学组件"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "多轴定制系统"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "支持多轴联动，满足复杂运动控制需求"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "嵌入式控制系统"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "集成化控制方案，开箱即用"
    p.level = 1
    
    # ==================== 第 5 页：业务与方案 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "业务与方案"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "🔧 非标自动化"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "节拍、工艺与安全联锁一次性想清楚，现场调试有清单可追溯"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "📊 数据看板"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "产量、OEE、停机原因分层呈现，支持大屏与移动端查看"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "📈 分期数字化"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "先做「看得见」再做「管得住」，阶段目标可对账、可验收"
    p.level = 1
    
    # ==================== 第 6 页：应用领域 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "应用领域"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    applications = [
        "半导体检测设备",
        "激光加工装备",
        "生物医疗仪器",
        "自动化产线",
        "高端装备制造"
    ]
    
    for app in applications:
        p = tf.add_paragraph()
        p.text = f"✓ {app}"
    
    # ==================== 第 7 页：合作与生态 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "合作与生态"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "我们与行业领军企业结成重要合作伙伴，共筑国产高端装备自主可控的坚实基石"
    
    partners = [
        "装备集团",
        "汽车零部件",
        "电子制造",
        "新材料",
        "医械配套",
        "包装机械"
    ]
    
    for partner in partners:
        p = tf.add_paragraph()
        p.text = f"• ×× {partner}"
    
    # ==================== 第 8 页：联系我们 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "联系我们"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "有项目想法？先从一次对齐开始"
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "告诉我们行业、现场痛点与期望周期"
    
    p = tf.add_paragraph()
    p.text = "我们会反馈可行路径与资源评估"
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "📍 地址：浙江温州"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🌐 网站：https://xdzn-tech.com/"
    
    p = tf.add_paragraph()
    p.text = "💼 合作理念：少承诺 · 多交付"
    
    # 保存文件
    output_path = "/home/tenbox/.openclaw/workspace/芯动智能科技公司介绍.pptx"
    prs.save(output_path)
    print(f"✅ PPT 生成成功：{output_path}")
    return output_path

if __name__ == "__main__":
    create_xdzn_ppt()
