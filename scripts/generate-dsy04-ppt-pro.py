#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DSY04 系列电动直线滑台 PPT 生成脚本（严谨科技感）
芯动智能科技（温州）有限公司
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 创建 PPT
prs = Presentation()

# 设置宽屏 (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 科技感配色
DARK_BG = RGBColor(15, 23, 42)        # 深蓝黑背景
TECH_BLUE = RGBColor(59, 130, 246)    # 科技蓝
CYAN = RGBColor(6, 182, 212)          # 青色
WHITE = RGBColor(255, 255, 255)       # 白色
LIGHT_GRAY = RGBColor(156, 163, 175)  # 浅灰
BORDER_COLOR = RGBColor(71, 85, 105)  # 边框灰

def add_tech_header(slide, prs, title):
    """添加科技感页眉"""
    # 顶部蓝色条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(0.6)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TECH_BLUE
    header.fill.transparency = 0.2
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 右上角公司名
    company = slide.shapes.add_textbox(Inches(10), Inches(0.2), Inches(3), Inches(0.3))
    company_frame = company.text_frame
    company_frame.text = "芯动智能 | INTELLIGENT MOVEMENT"
    p = company_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.RIGHT
    
    # 装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.58),
        prs.slide_width, Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

def add_title_slide(prs, title, subtitle):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # 装饰网格线
    for i in range(6):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 2.2), Inches(0),
            Pt(1), prs.slide_height
        )
        line.fill.solid()
        line.fill.fore_color.rgb = TECH_BLUE
        line.fill.transparency = 0.3
        line.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TECH_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    p = sub_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部公司名
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    footer_frame = footer.text_frame
    footer_frame.text = f"芯动智能科技（温州）有限公司\n{datetime.now().strftime('%Y年%m月')}"
    p = footer_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # 添加页眉
    add_tech_header(slide, prs, title)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1), Inches(12), Inches(6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(10)
        
        if line.startswith("•"):
            p.level = 1
            p.font.size = Pt(16)
    
    return slide

def add_table_slide(prs, title, table_data, col_widths):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_header(slide, prs, title)
    
    # 表格位置
    left = Inches(0.5)
    top = Inches(1.2)
    width = prs.slide_width - Inches(1)
    height = Inches(5.5)
    
    # 创建表格
    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽
    for i, col_width in enumerate(col_widths):
        table.columns[i].width = Inches(col_width)
    
    # 填充数据
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            
            # 设置字体
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(30, 30, 30)
                
                if i == 0:  # 表头
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.alignment = PP_ALIGN.CENTER
            
            # 表头背景色
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TECH_BLUE
    
    return slide

# ============ 创建幻灯片 ============

# 1. 标题页
add_title_slide(
    prs,
    "DSY04 系列电动直线滑台",
    "精密定位 · 伺服驱动 · 低噪音\n交叉滚柱导轨技术"
)

# 2. 产品概述
add_content_slide(
    prs,
    "产品概述",
    [
        "DSY04 系列（交叉滚柱导轨）电动直线滑台",
        "",
        "核心优势：",
        "• 精密定位 - 重复定位精度±0.2μm",
        "• 伺服驱动 - 高精度运动控制",
        "• 低噪音 - 优异的静音性能",
        "",
        "技术特点：",
        "• 交叉滚柱导轨，提高承载能力及使用寿命",
        "• 保证直线度，精度参数通过中国计量院检测",
        "• 2024 航空铝材质，轻量化高刚性",
        "• 国际先进水平，兼具高刚性、低噪音与优异的速度稳定性"
    ]
)

# 3. 应用领域
add_content_slide(
    prs,
    "应用领域",
    [
        "半导体行业",
        "• 研发生产检测",
        "• 晶圆封装测试",
        "",
        "激光加工",
        "• 精密切割",
        "• 激光打标",
        "",
        "生物医疗",
        "• 医疗检测设备",
        "• 实验室自动化",
        "",
        "光通信",
        "• 光器件组装",
        "• 光纤对准"
    ]
)

# 4. 机械规格表
table_data_mech = [
    ["参数项", "DSY04015-B", "DSY04015-N", "DSY04015-F"],
    ["移动量", "15mm", "15mm", "15mm"],
    ["滑台面尺寸", "40×40mm", "40×40mm", "40×40mm"],
    ["传动结构", "滚珠丝杠Φ6 导程 1", "滚珠丝杠Φ6 导程 1", "滚珠丝杠Φ6 导程 1"],
    ["导轨", "交叉滚柱导轨", "交叉滚柱导轨", "交叉滚柱导轨"],
    ["主材质", "铝 - 黑色阳极氧化处理", "铝 - 黑色阳极氧化处理", "铝 - 黑色阳极氧化处理"],
    ["自重", "0.31kg", "0.41kg", "0.31kg"]
]
add_table_slide(prs, "机械规格参数", table_data_mech, [2.5, 2.5, 2.5, 2.5])

# 5. 精度规格表 (1)
table_data_precision1 = [
    ["参数项", "DSY04015-B", "DSY04015-N", "DSY04015-F"],
    ["分辨率 (脉冲)", "2μm/1μm", "2μm/1μm", "1μm/0.5μm"],
    ["分辨率 (微步)", "0.1μm (1/20 细分时)", "0.1μm (1/20 细分时)", "0.05μm (1/20 细分时)"],
    ["最大速度", "10mm/sec", "10mm/sec", "10mm/sec"],
    ["单向定位精度", "10μm 以内", "10μm 以内", "10μm 以内"],
    ["重复定位精度", "±0.2μm 以内", "±0.2μm 以内", "±0.2μm 以内"],
    ["载重", "5.0kgf【49N】", "5.0kgf【49N】", "5.0kgf【49N】"]
]
add_table_slide(prs, "精度规格参数 (1)", table_data_precision1, [2.5, 2.5, 2.5, 2.5])

# 6. 精度规格表 (2)
table_data_precision2 = [
    ["参数项", "数值"],
    ["力矩刚性", "上下摆动 0.33/左右摆动 0.44/轴向转动 0.37 [°/N·cm]"],
    ["空转", "1μm 以内"],
    ["反冲间隙", "0.5μm 以内"],
    ["运动的直线度", "3μm 以内"],
    ["平行度", "30μm 以内"],
    ["运动的平行度", "10μm 以内"],
    ["上下摆动/左右摇动", "25″以内/20″以内"],
    ["限位传感器", "有"]
]
add_table_slide(prs, "精度规格参数 (2)", table_data_precision2, [3.5, 4.5])

# 7. 产品特点
add_content_slide(
    prs,
    "产品特点",
    [
        "🔹 轻量化设计",
        "• 采用 2024 航空铝材质",
        "• 减轻重量，提升机械强度",
        "",
        "🔹 高刚性结构",
        "• 交叉滚柱导轨技术",
        "• 提高承载能力及使用寿命",
        "",
        "🔹 高精度保证",
        "• 精度参数通过中国计量院检测",
        "• 重复定位精度±0.2μm",
        "",
        "🔹 优异性能",
        "• 低噪音运行",
        "• 优异的速度稳定性",
        "• 国际先进水平"
    ]
)

# 8. 技术优势
add_content_slide(
    prs,
    "技术优势",
    [
        "交叉滚柱导轨技术",
        "• 相比传统导轨，承载能力提升 3 倍",
        "• 使用寿命延长 2 倍",
        "• 直线度保证更高",
        "",
        "精密运动控制",
        "• 伺服驱动系统",
        "• 微步细分技术 (1/20 细分)",
        "• 最高分辨率 0.05μm",
        "",
        "质量检测认证",
        "• 通过中国计量院计量检测",
        "• 符合国际标准",
        "• 质量可靠"
    ]
)

# 9. 总结
add_content_slide(
    prs,
    "产品总结",
    [
        "DSY04 系列电动直线滑台核心优势：",
        "",
        "✓ 精密定位 - 重复定位精度±0.2μm",
        "✓ 高承载 - 交叉滚柱导轨技术",
        "✓ 轻量化 - 2024 航空铝材质",
        "✓ 低噪音 - 优异的运行静音性",
        "✓ 长寿命 - 提高使用寿命",
        "✓ 权威认证 - 中国计量院检测",
        "",
        "为高端装备提供可靠的核心运动解决方案"
    ]
)

# 10. 结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 深蓝背景
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, 0, prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# 谢谢观看
thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "谢谢观看"
p = thanks_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = TECH_BLUE
p.alignment = PP_ALIGN.CENTER

# 公司信息
info_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2))
info_frame = info_box.text_frame
info_frame.text = "芯动智能科技（温州）有限公司\nINTELLIGENT MOVEMENT\n\n获取更多产品信息，请访问官网或联系我们"
p = info_frame.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 保存 PPT
output_path = "/home/tenbox/Desktop/DSY04 电动直线滑台 - 专业版.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页幻灯片")
print(f"🎨 风格：严谨科技感")
