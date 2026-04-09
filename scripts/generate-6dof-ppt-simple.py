#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
六维耦合系统 PPT 生成脚本 (简化版) - 小爪 🐾
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# 创建 PPT
prs = Presentation()

def add_slide(prs, title, content=""):
    """添加幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # 内容
    if content:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content
        for p in tf.paragraphs:
            p.font.size = Pt(18)
    
    return slide

# 1. 标题页
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "六维耦合系统"
slide.placeholders[1].text = "Six-Degree-of-Freedom Coupled System\n技术原理与应用\n" + datetime.now().strftime("%Y年%m月")

# 2. 目录
add_slide(prs, "目录", 
    "1. 六维系统概述\n"
    "2. 耦合系统原理\n"
    "3. 数学模型\n"
    "4. 关键技术\n"
    "5. 应用领域\n"
    "6. 发展趋势"
)

# 3. 六维系统概述
add_slide(prs, "六维系统概述",
    "什么是六维系统？\n"
    "• 六自由度（6-DOF）运动系统\n"
    "• 三个平移自由度：X、Y、Z 轴\n"
    "• 三个旋转自由度：Roll、Pitch、Yaw\n\n"
    "系统特点：\n"
    "• 全空间运动能力\n"
    "• 高精度定位控制\n"
    "• 多轴协同耦合"
)

# 4. 耦合系统原理
add_slide(prs, "耦合系统原理",
    "耦合的定义\n"
    "• 多个子系统相互影响、相互作用\n"
    "• 输入输出之间存在关联关系\n\n"
    "耦合类型：\n"
    "• 机械耦合：结构连接导致的力/力矩传递\n"
    "• 控制耦合：多控制器之间的信号交互\n"
    "• 动力耦合：能量在系统间的传递"
)

# 5. 数学模型
add_slide(prs, "数学模型",
    "运动学方程：\n"
    "• 位置向量：p = [x, y, z]ᵀ\n"
    "• 姿态向量：θ = [α, β, γ]ᵀ\n\n"
    "动力学方程：\n"
    "• M(q)q̈ + C(q,q̇)q̇ + G(q) = τ\n"
    "• M: 质量矩阵\n"
    "• C: 科里奥利力矩阵\n"
    "• G: 重力向量\n"
    "• τ: 广义力/力矩"
)

# 6. 关键技术
add_slide(prs, "关键技术",
    "1. 精密驱动技术\n"
    "• 伺服电机/直线电机\n"
    "• 压电陶瓷驱动器\n\n"
    "2. 传感检测技术\n"
    "• 激光干涉仪\n"
    "• 光电编码器\n"
    "• 惯性测量单元 (IMU)\n\n"
    "3. 控制算法\n"
    "• PID 控制\n"
    "• 自适应控制\n"
    "• 鲁棒控制"
)

# 7. 应用领域
add_slide(prs, "应用领域",
    "🔬 科学实验\n"
    "• 振动模拟测试台\n"
    "• 微重力环境模拟\n\n"
    "🏭 工业制造\n"
    "• 精密加工定位\n"
    "• 机器人末端执行器\n\n"
    "✈️ 航空航天\n"
    "• 飞行模拟器\n"
    "• 卫星姿态控制\n\n"
    "🏥 医疗设备\n"
    "• 手术机器人\n"
    "• 康复训练设备"
)

# 8. 发展趋势
add_slide(prs, "发展趋势",
    "1. 高精度化\n"
    "• 纳米级定位精度\n\n"
    "2. 高速化\n"
    "• 高频响驱动\n\n"
    "3. 智能化\n"
    "• AI 辅助控制\n\n"
    "4. 集成化\n"
    "• 模块化设计"
)

# 9. 总结
add_slide(prs, "总结",
    "六维耦合系统核心价值：\n"
    "• 实现全空间六自由度精确控制\n"
    "• 解决多轴协同耦合难题\n"
    "• 支撑高端装备核心技术\n\n"
    "未来展望：\n"
    "• 更智能、更精密、更集成"
)

# 10. 结束页
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "谢谢观看"
slide.placeholders[1].text = "Q & A\n\n小爪 🐾 自动生成"

# 保存 PPT
output_path = "/home/tenbox/Desktop/六维耦合系统.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
