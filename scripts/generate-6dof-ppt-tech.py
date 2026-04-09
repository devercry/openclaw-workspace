#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
六维耦合系统 PPT 生成脚本 (科技感精美版) - 小爪 🐾
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 创建 PPT
prs = Presentation()

# 设置宽屏 (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 科技感配色
DARK_BLUE = RGBColor(10, 25, 47)      # 深蓝背景
CYAN = RGBColor(0, 210, 255)          # 青色高亮
BLUE = RGBColor(0, 150, 255)          # 蓝色
PURPLE = RGBColor(138, 43, 226)       # 紫色
WHITE = RGBColor(255, 255, 255)       # 白色
LIGHT_GRAY = RGBColor(200, 200, 200)  # 浅灰

def add_tech_background(slide, prs):
    """添加科技感背景"""
    # 深蓝背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # 装饰线条
    for i in range(5):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 2.5), Inches(7), Inches(3), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = CYAN
        line.fill.transparency = 0.3
        line.line.fill.background()

def add_title_slide(prs, title, subtitle):
    """添加标题页（科技感）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    add_tech_background(slide, prs)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    p = sub_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 装饰六边形
    hexagon = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(5.5), Inches(5.5),
        Inches(2.5), Inches(2)
    )
    hexagon.fill.solid()
    hexagon.fill.fore_color.rgb = BLUE
    hexagon.fill.transparency = 0.7
    hexagon.line.color.rgb = CYAN
    hexagon.line.width = Pt(2)
    
    return slide

def add_content_slide(prs, title, items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_background(slide, prs)
    
    # 标题栏背景
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.fill.transparency = 0.3
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        if item.startswith("•"):
            p.level = 1
            p.font.color.rgb = LIGHT_GRAY
    
    return slide

def add_diagram_slide(prs, title, diagram_desc):
    """添加图示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_tech_background(slide, prs)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # 中心六边形
    hex_center = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(5.5), Inches(2.5),
        Inches(3), Inches(2.5)
    )
    hex_center.fill.solid()
    hex_center.fill.fore_color.rgb = BLUE
    hex_center.fill.transparency = 0.6
    hex_center.line.color.rgb = CYAN
    hex_center.line.width = Pt(3)
    
    # 六个顶点小圆
    positions = [
        (Inches(5.5), Inches(1.2)),   # 上
        (Inches(7.5), Inches(2)),     # 右上
        (Inches(7.5), Inches(4.5)),   # 右下
        (Inches(5.5), Inches(5.3)),   # 下
        (Inches(3.5), Inches(4.5)),   # 左下
        (Inches(3.5), Inches(2)),     # 左上
    ]
    labels = ["X", "Y", "Z", "Roll", "Pitch", "Yaw"]
    
    for i, (x, y) in enumerate(positions):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - Inches(0.3), y - Inches(0.3),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.color.rgb = CYAN
        circle.line.width = Pt(2)
        
        # 标签
        label = slide.shapes.add_textbox(x - Inches(0.5), y - Inches(0.8), Inches(1), Inches(0.5))
        label_frame = label.text_frame
        label_frame.text = labels[i]
        label_frame.paragraphs[0].font.size = Pt(16)
        label_frame.paragraphs[0].font.color.rgb = WHITE
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 说明
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = diagram_desc
    desc_frame.paragraphs[0].font.size = Pt(18)
    desc_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 创建幻灯片 ============

# 1. 标题页
add_title_slide(
    prs,
    "六维耦合系统",
    "Six-Degree-of-Freedom Coupled System\n技术原理与应用\n" + datetime.now().strftime("%Y年%m月")
)

# 2. 目录
add_content_slide(
    prs,
    "目录 CONTENTS",
    [
        "01  六维系统概述",
        "02  耦合系统原理",
        "03  数学模型",
        "04  关键技术",
        "05  应用领域",
        "06  发展趋势"
    ]
)

# 3. 六维系统概述
add_content_slide(
    prs,
    "六维系统概述",
    [
        "什么是六维系统？",
        "• 六自由度（6-DOF）运动系统",
        "• 三个平移自由度：X、Y、Z 轴",
        "• 三个旋转自由度：Roll、Pitch、Yaw",
        "",
        "系统特点：",
        "• 全空间运动能力",
        "• 高精度定位控制",
        "• 多轴协同耦合"
    ]
)

# 4. 六自由度示意图
add_diagram_slide(
    prs,
    "六自由度示意图",
    "六个自由度：X/Y/Z 平移 + Roll/Pitch/Yaw 旋转，实现全空间精确控制"
)

# 5. 耦合系统原理
add_content_slide(
    prs,
    "耦合系统原理",
    [
        "耦合的定义",
        "• 多个子系统相互影响、相互作用",
        "• 输入输出之间存在关联关系",
        "",
        "耦合类型：",
        "• 机械耦合：结构连接导致的力/力矩传递",
        "• 控制耦合：多控制器之间的信号交互",
        "• 动力耦合：能量在系统间的传递",
        "",
        "解耦控制：将多变量系统分解为单变量系统"
    ]
)

# 6. 数学模型
add_content_slide(
    prs,
    "数学模型",
    [
        "运动学方程：",
        "• 位置向量：p = [x, y, z]ᵀ",
        "• 姿态向量：θ = [α, β, γ]ᵀ",
        "• 广义坐标：q = [pᵀ, θᵀ]ᵀ",
        "",
        "动力学方程：",
        "• M(q)q̈ + C(q,q̇)q̇ + G(q) = τ",
        "• M: 质量矩阵",
        "• C: 科里奥利力矩阵",
        "• G: 重力向量",
        "• τ: 广义力/力矩"
    ]
)

# 7. 关键技术
add_content_slide(
    prs,
    "关键技术",
    [
        "🔹 精密驱动技术",
        "• 伺服电机/直线电机",
        "• 压电陶瓷驱动器",
        "",
        "🔹 传感检测技术",
        "• 激光干涉仪",
        "• 光电编码器",
        "• 惯性测量单元 (IMU)",
        "",
        "🔹 控制算法",
        "• PID 控制",
        "• 自适应控制",
        "• 鲁棒控制",
        "• 智能控制"
    ]
)

# 8. 应用领域
add_content_slide(
    prs,
    "应用领域",
    [
        "🔬 科学实验",
        "• 振动模拟测试台",
        "• 微重力环境模拟",
        "",
        "🏭 工业制造",
        "• 精密加工定位",
        "• 机器人末端执行器",
        "",
        "✈️ 航空航天",
        "• 飞行模拟器",
        "• 卫星姿态控制",
        "",
        "🏥 医疗设备",
        "• 手术机器人",
        "• 康复训练设备"
    ]
)

# 9. 发展趋势
add_content_slide(
    prs,
    "发展趋势",
    [
        "📈 高精度化",
        "• 纳米级定位精度",
        "• 亚微弧度角精度",
        "",
        "⚡ 高速化",
        "• 高频响驱动",
        "• 实时控制优化",
        "",
        "🤖 智能化",
        "• AI 辅助控制",
        "• 自学习自适应",
        "",
        "🔗 集成化",
        "• 模块化设计",
        "• 系统小型化"
    ]
)

# 10. 总结
add_content_slide(
    prs,
    "总结",
    [
        "💎 六维耦合系统核心价值：",
        "• 实现全空间六自由度精确控制",
        "• 解决多轴协同耦合难题",
        "• 支撑高端装备核心技术",
        "",
        "🚀 未来展望：",
        "• 更智能、更精密、更集成"
    ]
)

# 11. 结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_tech_background(slide, prs)

# 谢谢观看
thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "谢谢观看"
p = thanks_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

# Q&A
qa_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1.5))
qa_frame = qa_box.text_frame
qa_frame.text = "Q & A\n\n小爪 🐾 自动生成"
p = qa_frame.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 保存 PPT
output_path = "/home/tenbox/Desktop/六维耦合系统 - 科技感.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页幻灯片")
print(f"🎨 科技感配色：深蓝背景 + 青色高亮")
